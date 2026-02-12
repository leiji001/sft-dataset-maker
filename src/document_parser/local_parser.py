"""本地文档解析器 - 支持 PDF、DOCX、TXT、PPTX"""

from pathlib import Path

import PyPDF2
import docx
from pptx import Presentation


def has_images_in_pdf(file_path: Path) -> bool:
    """检测 PDF 文件是否包含图片"""
    try:
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                resources = page.get("/Resources")
                if resources and "/XObject" in resources:
                    xobjects = resources["/XObject"].get_object()
                    for obj_name in xobjects:
                        xobj = xobjects[obj_name].get_object()
                        if xobj.get("/Subtype") == "/Image":
                            return True
    except Exception:
        pass
    return False


def parse_pdf(file_path: Path) -> str:
    """解析 PDF 文件的纯文本内容"""
    text_parts: list[str] = []
    with open(file_path, "rb") as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text.strip())
    return "\n\n".join(text_parts)


def parse_docx(file_path: Path) -> str:
    """解析 DOCX 文件"""
    doc = docx.Document(str(file_path))
    text_parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text.strip())
    # 提取表格文本
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text_parts.append(row_text)
    return "\n\n".join(text_parts)


def has_images_in_docx(file_path: Path) -> bool:
    """检测 DOCX 文件是否包含图片"""
    try:
        doc = docx.Document(str(file_path))
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                return True
    except Exception:
        pass
    return False


def parse_pptx(file_path: Path) -> str:
    """解析 PPTX 文件"""
    prs = Presentation(str(file_path))
    text_parts: list[str] = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text.strip() for cell in row.cells if cell.text.strip()
                    )
                    if row_text:
                        slide_texts.append(row_text)
        if slide_texts:
            text_parts.append(f"[幻灯片 {slide_idx}]\n" + "\n".join(slide_texts))
    return "\n\n".join(text_parts)


def has_images_in_pptx(file_path: Path) -> bool:
    """检测 PPTX 文件是否包含图片"""
    try:
        prs = Presentation(str(file_path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    return True
    except Exception:
        pass
    return False


def parse_txt(file_path: Path) -> str:
    """解析纯文本文件"""
    encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
    for enc in encodings:
        try:
            return file_path.read_text(encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    raise ValueError(f"无法解码文件: {file_path}")


def parse_markdown(file_path: Path) -> str:
    """解析 Markdown 文件"""
    return parse_txt(file_path)


# 支持的文件格式与对应的解析函数
PARSERS: dict[str, callable] = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".doc": parse_docx,
    ".pptx": parse_pptx,
    ".txt": parse_txt,
    ".md": parse_markdown,
}

# 支持的图片检测函数
IMAGE_CHECKERS: dict[str, callable] = {
    ".pdf": has_images_in_pdf,
    ".docx": has_images_in_docx,
    ".pptx": has_images_in_pptx,
}
